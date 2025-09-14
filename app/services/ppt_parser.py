from pptx import Presentation

def parse_ppt(file_path: str) -> dict:
    prs = Presentation(file_path)
    slides_text = {}

    for i, slide in enumerate(prs.slides, start=1):
        texts = []

        for shape in slide.shapes:
            # Shape에 text 속성 있으면 리스트에 추가
            if hasattr(shape, "text"):
                texts.append(shape.text.strip())
                
        slides_text[f"slide_{i}"] = " ".join(texts)

    return slides_text
