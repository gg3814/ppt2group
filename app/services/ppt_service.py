from pptx import Presentation
import logging

log = logging.getLogger(__name__)

# {슬라이드번호: "텍스트1 텍스트2 텍스트3"} 형식의 딕셔너리 반환
def parse_ppt(file_path: str) -> dict[int, str]:
    prs = Presentation(file_path)
    slides_text = {}

    for i, slide in enumerate(prs.slides, start=1):
        texts = parse_texts(slide)
        slides_text[i] = " ".join(texts)

    log.info("Parsed file's path: %s", file_path)
    return slides_text

# slide: Slide 타입
def parse_texts(slide) -> list[str]:
    texts = []

    for shape in slide.shapes:
        if hasattr(shape, "text"):
            texts.append(shape.text.strip())
    
    return texts

